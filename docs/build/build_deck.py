"""
Builds docs/NeurAI-Platform-Demo.pptx — the demo deck.

Fourteen slides: what the platform is, why it exists, what Echo does, how it
is built, and where it is heading. Widescreen 16:9, the product's own palette.

Run:  python docs/build/build_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt

ROOT = Path(__file__).resolve().parents[2]
DIAGRAMS = ROOT / "docs" / "diagrams" / "png"
SHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "NeurAI-Platform-Demo.pptx"

W, H = Cm(33.867), Cm(19.05)          # 16:9
M = Cm(2.0)                            # the page margin

INK = RGBColor(0x16, 0x12, 0x1F)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF7, 0xF4, 0xFC)
MUTED = RGBColor(0x6B, 0x64, 0x78)
LINE = RGBColor(0xDA, 0xD4, 0xE4)
VIOLET = RGBColor(0x7C, 0x4D, 0xFF)
VIOLET_D = RGBColor(0x5B, 0x34, 0xC7)
VIOLET_PALE = RGBColor(0xEF, 0xE8, 0xFF)
PINK = RGBColor(0xE0, 0x5A, 0x9B)
PINK_PALE = RGBColor(0xFD, 0xEA, 0xF3)
BLUE = RGBColor(0x44, 0x68, 0xD8)
BLUE_PALE = RGBColor(0xE8, 0xED, 0xFC)
GREEN = RGBColor(0x1F, 0x8A, 0x5B)
GREEN_PALE = RGBColor(0xE4, 0xF5, 0xEC)
AMBER = RGBColor(0xB4, 0x70, 0x0E)
AMBER_PALE = RGBColor(0xFC, 0xF1, 0xDF)
INK_DARK = RGBColor(0x0C, 0x09, 0x22)

FONT = "Segoe UI"
MONO = "Consolas"


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs, fill=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs = [(text, size, bold, colour, font?)] — one paragraph per entry
    unless the text is '' which inserts a spacer."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in runs:
        content, size, bold, colour = item[0], item[1], item[2], item[3]
        font = item[4] if len(item) > 4 else FONT
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = content
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = font
    return box


def card(slide, left, top, width, height, fill, *, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.06
        except (IndexError, KeyError):
            pass
    return shape


def eyebrow(slide, label):
    text(slide, M, Cm(1.35), W - 2 * M, Cm(0.8),
         [(label.upper(), 10, True, VIOLET)])


def title(slide, heading, sub=None):
    text(slide, M, Cm(2.15), W - 2 * M, Cm(2.0),
         [(heading, 30, True, INK)])
    if sub:
        text(slide, M, Cm(3.85), Cm(24), Cm(1.6),
             [(sub, 13, False, MUTED)], spacing=1.15)


def footer(slide, n):
    text(slide, M, H - Cm(1.3), Cm(16), Cm(0.7),
         [("NeurAI Platform", 8.5, False, MUTED)])
    text(slide, W - M - Cm(3), H - Cm(1.3), Cm(3), Cm(0.7),
         [(str(n), 8.5, False, MUTED)], align=PP_ALIGN.RIGHT)


def picture(slide, path, left, top, width):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width)


# ------------------------------------------------------------------ slides

def slide_cover(prs):
    s = blank(prs, INK_DARK)
    card(s, Cm(0), Cm(0), Cm(14), H, RGBColor(0x19, 0x10, 0x34), radius=False)
    text(s, M, Cm(6.4), Cm(20), Cm(3),
         [("NeurAI", 54, True, RGBColor(0xC4, 0xB0, 0xFF))])
    text(s, M, Cm(8.9), Cm(24), Cm(2),
         [("The meeting is the record.", 24, False, PAPER)])
    text(s, M, Cm(11.1), Cm(21), Cm(3),
         [("A Persian-first AI platform for organisations —\n"
           "and Echo, the application that turns every\n"
           "conversation into something you can search,\n"
           "act on, and prove.", 13, False, RGBColor(0xA9, 0xA0, 0xC0))],
         spacing=1.3)
    text(s, M, H - Cm(2.4), Cm(20), Cm(1),
         [("Product demonstration · August 2026", 10, False, RGBColor(0x7A, 0x71, 0x94))])
    return s


def slide_problem(prs):
    s = blank(prs)
    eyebrow(s, "The problem")
    title(s, "Everything important is said out loud,\nand then it is gone.")
    items = [
        ("What was decided?", "Nobody wrote it down. Two people remember it differently.",
         PINK_PALE, PINK),
        ("Who committed to what?", "It was clear in the room. It is not clear on Thursday.",
         AMBER_PALE, AMBER),
        ("Has this come up before?", "Probably. There is no way to look.",
         BLUE_PALE, BLUE),
    ]
    x = M
    width = Cm(9.3)
    for head, sub, fill, accent in items:
        card(s, x, Cm(7.6), width, Cm(6.2), fill)
        card(s, x + Cm(0.9), Cm(8.5), Cm(0.5), Cm(0.5), accent)
        text(s, x + Cm(0.9), Cm(9.5), width - Cm(1.8), Cm(1.4),
             [(head, 17, True, INK)])
        text(s, x + Cm(0.9), Cm(11.1), width - Cm(1.8), Cm(2.4),
             [(sub, 12, False, MUTED)], spacing=1.25)
        x += width + Cm(1.35)
    text(s, M, Cm(15.2), W - 2 * M, Cm(1.4),
         [("Meetings are where organisations actually decide things — and the least "
           "recorded part of the business.", 13, True, INK)])
    footer(s, 2)
    return s


def slide_what(prs):
    s = blank(prs)
    eyebrow(s, "What we built")
    title(s, "A platform, and its first application.")

    card(s, M, Cm(6.4), Cm(14.4), Cm(9.6), VIOLET)
    text(s, M + Cm(1.2), Cm(7.4), Cm(12), Cm(1.2),
         [("NeurAI Platform", 22, True, PAPER)])
    text(s, M + Cm(1.2), Cm(9.0), Cm(12), Cm(6),
         [("The shell every application lives in.", 13, True, PAPER),
          ("Identity, organisations, roles and permissions.", 12, False, RGBColor(0xE4, 0xDC, 0xFF)),
          ("An assistant that reaches across everything you are allowed to see.", 12, False, RGBColor(0xE4, 0xDC, 0xFF)),
          ("Persian-first: right-to-left, Jalali dates, Persian digits.", 12, False, RGBColor(0xE4, 0xDC, 0xFF))],
         spacing=1.25)

    card(s, M + Cm(15.4), Cm(6.4), Cm(14.4), Cm(9.6), WASH, line=LINE)
    text(s, M + Cm(16.6), Cm(7.4), Cm(12), Cm(1.2),
         [("Echo", 22, True, VIOLET_D)])
    text(s, M + Cm(16.6), Cm(9.0), Cm(12), Cm(6),
         [("The first application: calls and meetings.", 13, True, INK),
          ("Record in the browser, or upload what you already have.", 12, False, MUTED),
          ("Persian transcription at 2.1% word error rate.", 12, False, MUTED),
          ("Versioned summaries with actions, decisions and named speakers.", 12, False, MUTED)],
         spacing=1.25)
    footer(s, 3)
    return s


def slide_why_platform(prs):
    s = blank(prs)
    eyebrow(s, "Why a platform, not a product")
    title(s, "The hard parts are the shared parts.",
          "Transcription is a purchased capability. What takes years is everything "
          "around it — and it is the same everything for the next application.")
    rows = [
        ("Who can see what", "An organisation's meetings are its most sensitive data. "
         "Permissions are enforced in the database, not in the interface."),
        ("An assistant with limits", "It answers with exactly the caller's own access, "
         "and it cannot change anything without a person confirming it."),
        ("Persian done properly", "Not a translated interface. Direction, digits, "
         "calendars and text normalisation are structural."),
        ("A record you can defend", "Versioned artifacts, provenance on everything, "
         "and an audit trail written in the same transaction as the change."),
    ]
    y = Cm(7.4)
    for head, sub in rows:
        card(s, M, y, W - 2 * M, Cm(2.1), WASH, line=LINE)
        text(s, M + Cm(0.9), y + Cm(0.42), Cm(8), Cm(1.2),
             [(head, 14, True, VIOLET_D)])
        text(s, M + Cm(9.6), y + Cm(0.5), Cm(19), Cm(1.4),
             [(sub, 11.5, False, INK)], spacing=1.15)
        y += Cm(2.5)
    footer(s, 4)
    return s


def slide_flow(prs):
    s = blank(prs)
    eyebrow(s, "How Echo works")
    title(s, "Four steps, and nothing to remember.")
    steps = [
        ("1", "Record", "In the browser, or upload a file you already have. "
         "Thirty-minute parts, so nothing is lost."),
        ("2", "Transcribe", "Silence is removed first, then Persian speech becomes "
         "text with word-level timing."),
        ("3", "Attribute", "Voices are separated and matched to people you have "
         "enrolled. You confirm the rest."),
        ("4", "Summarise", "A summary against the template you chose — actions, "
         "decisions, tags. Versioned, never overwritten."),
    ]
    x = M
    width = Cm(6.9)
    for num, head, sub in steps:
        card(s, x, Cm(7.2), width, Cm(7.4), PAPER, line=LINE)
        circle = card(s, x + Cm(0.85), Cm(8.0), Cm(1.25), Cm(1.25), VIOLET)
        text(s, x + Cm(0.85), Cm(8.22), Cm(1.25), Cm(0.9),
             [(num, 15, True, PAPER)], align=PP_ALIGN.CENTER)
        text(s, x + Cm(0.85), Cm(9.9), width - Cm(1.7), Cm(1),
             [(head, 17, True, INK)])
        text(s, x + Cm(0.85), Cm(11.2), width - Cm(1.7), Cm(3),
             [(sub, 11.5, False, MUTED)], spacing=1.25)
        x += width + Cm(0.72)
    text(s, M, Cm(15.6), W - 2 * M, Cm(1.2),
         [("Every step is a queue message. A part that fails leaves a visible gap; "
           "the rest of the meeting still completes.", 12, False, INK)])
    footer(s, 5)
    return s


def slide_pipeline_figure(prs):
    s = blank(prs)
    eyebrow(s, "The pipeline")
    title(s, "And when a step cannot do its job.")
    picture(s, DIAGRAMS / "deck-pipeline.png", M + Cm(1.6), Cm(5.2), Cm(26.6))
    footer(s, 6)
    return s


def slide_record(prs):
    s = blank(prs)
    eyebrow(s, "The record")
    title(s, "One page holds the whole meeting.")
    picture(s, SHOTS / "record-page.png", M, Cm(5.6), Cm(20.2))
    notes = [
        ("The summary, versioned", "Regenerate against a different template and the "
         "old version stays."),
        ("The transcript beneath it", "Click any line to hear it. One timeline across "
         "every part."),
        ("Speakers you can name", "Link a voice to a person once and it is recognised "
         "in future meetings."),
        ("Corrections are kept", "An edited line is marked as human-edited. The record "
         "says who changed it."),
    ]
    y = Cm(6.0)
    for head, sub in notes:
        text(s, M + Cm(21), y, Cm(8.6), Cm(0.8), [(head, 12.5, True, VIOLET_D)])
        text(s, M + Cm(21), y + Cm(0.72), Cm(8.6), Cm(1.8),
             [(sub, 10.5, False, MUTED)], spacing=1.2)
        y += Cm(2.5)
    footer(s, 7)
    return s


def slide_assistant(prs):
    s = blank(prs)
    eyebrow(s, "The assistant")
    title(s, "It answers with your access, and it asks before it acts.")
    left = [
        ("Ask across every meeting", "\"What did we decide about the budget?\" — it "
         "searches the transcripts you are allowed to read, and cites the meeting."),
        ("It cannot see more than you", "Every tool it calls runs under your identity. "
         "There is no privileged path and no admin mode."),
    ]
    right = [
        ("Changes are proposals", "It never writes. It describes the change, shows you "
         "the before and the after, and waits."),
        ("Your answer is the record", "Approve or reject — both are recorded as a human "
         "decision, and the assistant cannot read them back."),
    ]
    y0 = Cm(7.0)
    for col, items in ((M, left), (M + Cm(15.4), right)):
        y = y0
        for head, sub in items:
            card(s, col, y, Cm(14.4), Cm(3.6), WASH, line=LINE)
            text(s, col + Cm(0.9), y + Cm(0.5), Cm(12.6), Cm(0.9),
                 [(head, 14, True, VIOLET_D)])
            text(s, col + Cm(0.9), y + Cm(1.45), Cm(12.6), Cm(2),
                 [(sub, 11.5, False, INK)], spacing=1.2)
            y += Cm(4.2)
    card(s, M, Cm(15.6), W - 2 * M, Cm(1.7), VIOLET_PALE)
    text(s, M + Cm(0.9), Cm(15.95), W - 2 * M - Cm(1.8), Cm(1),
         [("The limit is enforced in the database, not in the prompt. "
           "A prompt is a request; a grant is a wall.", 12.5, True, VIOLET_D)])
    footer(s, 8)
    return s


def slide_dashboard(prs):
    s = blank(prs)
    eyebrow(s, "The landing page")
    title(s, "A board you arrange, not a report you receive.")
    text(s, M, Cm(5.6), Cm(13), Cm(8),
         [("Fifteen widgets, four sizes each.", 14, True, INK),
          ("Drag them where you want them. A tile you leave at the bottom stays at the "
           "bottom — the board does not rearrange itself under you.", 11.5, False, MUTED),
          ("", 6, False, MUTED),
          ("A bigger tile says more.", 14, True, INK),
          ("Each size is designed, not stretched. Growing a card adds rows, range or a "
           "chart — it never scales the same content up.", 11.5, False, MUTED),
          ("", 6, False, MUTED),
          ("It never invents a number.", 14, True, INK),
          ("Not-read-yet, could-not-read and genuinely-empty are three different states, "
           "and the board says which one it is in.", 11.5, False, MUTED)],
         spacing=1.25)
    picture(s, SHOTS / "hub-fa.png", M + Cm(14.6), Cm(5.9), Cm(15.2))
    footer(s, 9)
    return s


def slide_architecture(prs):
    s = blank(prs)
    eyebrow(s, "How it is built")
    title(s, "Five processes, one wall.")
    picture(s, DIAGRAMS / "deck-topology.png", M + Cm(1.2), Cm(5.3), Cm(27.4))
    text(s, M, Cm(15.9), W - 2 * M, Cm(1.6),
         [("The browser never holds a credential and never learns the API's address. "
           "The database decides what every identity may see, so the guarantee holds "
           "even if every layer above it is wrong.", 12, False, INK)], spacing=1.2)
    footer(s, 10)
    return s


def slide_security(prs):
    s = blank(prs)
    eyebrow(s, "What it refuses to do")
    title(s, "The absences are the feature.")
    items = [
        ("No content in logs", "Transcripts, summaries and messages are never logged. "
         "Database errors are logged by code, never by message.", GREEN_PALE, GREEN),
        ("No deletion by the application", "The product cannot physically delete a row. "
         "One separate process holds that right and runs at 03:30.", VIOLET_PALE, VIOLET_D),
        ("No fabricated numbers", "A value we could not read renders as an unknown, in "
           "words — never as a zero.", AMBER_PALE, AMBER),
        ("No arbitrary outbound requests", "Webhook addresses are checked as the socket "
         "opens, on the address actually being reached.", BLUE_PALE, BLUE),
        ("No token in the browser", "An injected script cannot read a credential that "
         "is not there.", PINK_PALE, PINK),
        ("No silent degradation", "Whatever is forfeited is said out loud, on the "
         "screen where it matters.", WASH, MUTED),
    ]
    x, y = M, Cm(6.8)
    width, height = Cm(9.3), Cm(4.4)
    for i, (head, sub, fill, accent) in enumerate(items):
        card(s, x, y, width, height, fill)
        text(s, x + Cm(0.8), y + Cm(0.55), width - Cm(1.6), Cm(1),
             [(head, 13.5, True, accent)])
        text(s, x + Cm(0.8), y + Cm(1.7), width - Cm(1.6), Cm(2.4),
             [(sub, 11, False, INK)], spacing=1.2)
        if i % 3 == 2:
            x = M
            y += height + Cm(0.9)
        else:
            x += width + Cm(1.35)
    footer(s, 11)
    return s


def slide_proof(prs):
    s = blank(prs)
    eyebrow(s, "Where it stands")
    title(s, "Built, deployed, and in use.")
    stats = [
        ("2.1%", "Persian word error rate", "Measured on a real recording, after "
         "normalisation."),
        ("98", "database migrations", "Hand-written SQL. Every one checksummed and "
         "append-only."),
        ("32", "tables, all walled", "Row-level security on every one, with 86 policies."),
        ("1,400+", "automated checks", "Across four suites, plus a real production build "
         "before anything ships."),
    ]
    x = M
    width = Cm(6.9)
    for big, label, sub in stats:
        card(s, x, Cm(6.8), width, Cm(6.4), WASH, line=LINE)
        text(s, x + Cm(0.8), Cm(7.5), width - Cm(1.6), Cm(2),
             [(big, 34, True, VIOLET_D)])
        text(s, x + Cm(0.8), Cm(9.9), width - Cm(1.6), Cm(0.9),
             [(label, 12, True, INK)])
        text(s, x + Cm(0.8), Cm(10.9), width - Cm(1.6), Cm(2),
             [(sub, 10.5, False, MUTED)], spacing=1.2)
        x += width + Cm(0.72)
    card(s, M, Cm(14.2), W - 2 * M, Cm(2.6), VIOLET_PALE)
    text(s, M + Cm(1.0), Cm(14.75), W - 2 * M - Cm(2), Cm(1.8),
         [("Running in production today: the web application on a global edge, the API "
           "and workers on a dedicated host behind a tunnel, and a managed Postgres that "
           "enforces every permission itself.", 12.5, False, INK)], spacing=1.2)
    footer(s, 12)
    return s


def slide_next(prs):
    s = blank(prs)
    eyebrow(s, "Where we are heading")
    title(s, "From a record of what happened\nto an assistant that acts on it.")
    lanes = [
        ("Now", "Echo, complete", [
            "Record, transcribe, attribute, summarise",
            "Search across everything you may see",
            "The assistant answers and proposes",
            "Organisations, roles, audit, integrations",
        ], VIOLET, PAPER),
        ("Next", "The assistant acts", [
            "Open questions nobody answered",
            "Decisions that reverse earlier ones",
            "Preparation before the meeting starts",
            "Terms you choose to watch for",
        ], VIOLET_PALE, INK),
        ("Then", "More applications", [
            "The platform's shell is application-agnostic",
            "Identity, permissions and the assistant are shared",
            "Each new application inherits the wall",
            "Persian-first from the first line",
        ], WASH, INK),
    ]
    x = M
    width = Cm(9.3)
    for when, head, points, fill, ink in lanes:
        card(s, x, Cm(7.0), width, Cm(8.4), fill,
             line=None if fill != WASH else LINE)
        text(s, x + Cm(0.9), Cm(7.6), width - Cm(1.8), Cm(0.8),
             [(when.upper(), 10, True, ink if fill == VIOLET else VIOLET_D)])
        text(s, x + Cm(0.9), Cm(8.5), width - Cm(1.8), Cm(1.2),
             [(head, 17, True, ink)])
        runs = [("· " + p, 11.5, False, ink) for p in points]
        text(s, x + Cm(0.9), Cm(10.2), width - Cm(1.8), Cm(5), runs, spacing=1.3)
        x += width + Cm(1.35)
    footer(s, 13)
    return s


def slide_close(prs):
    s = blank(prs, INK_DARK)
    text(s, M, Cm(6.6), Cm(24), Cm(3),
         [("The meeting is the record.", 36, True, PAPER)])
    text(s, M, Cm(9.6), Cm(22), Cm(3),
         [("Everything said, kept — searchable, attributable,\n"
           "and never leaving the walls of your organisation.", 15, False,
           RGBColor(0xA9, 0xA0, 0xC0))], spacing=1.3)
    card(s, M, Cm(13.4), Cm(9), Cm(0.12), VIOLET, radius=False)
    text(s, M, Cm(14.4), Cm(20), Cm(1.4),
         [("neurai.pt", 16, True, RGBColor(0xC4, 0xB0, 0xFF))])
    return s


def build() -> Path:
    prs = deck()
    slide_cover(prs)
    slide_problem(prs)
    slide_what(prs)
    slide_why_platform(prs)
    slide_flow(prs)
    slide_pipeline_figure(prs)
    slide_record(prs)
    slide_assistant(prs)
    slide_dashboard(prs)
    slide_architecture(prs)
    slide_security(prs)
    slide_proof(prs)
    slide_next(prs)
    slide_close(prs)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print("wrote", build())
